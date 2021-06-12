import speechUtil as speechUtil
import os
import glob
import PyPDF2
import docx2txt
from pptx import Presentation
from nlp.lex.lexical_analyzer import LexicalAnalyzer

def readFromFile(fileName):
    filePath = "uploads/" + fileName
    filetype = fileName.split('.')[-1]
    fileNameWithoutExt = ( fileName.rsplit( ".", 1 )[ 0 ] )
    success = True
    text = ''
    
    try:
        if filetype == 'txt':
            if is_file_empty(filePath):
                success = False
                raise Exception("File string is empty!")
            else:
                with open(filePath,"r") as f:
                    text = f.read()
                    speechUtil.synthesize_and_save_to_file(text, fileNameWithoutExt)
        elif filetype == 'pdf':
            pdfFileObj = open(filePath, "rb")
            pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
            totalPages = pdfReader.numPages
            for i in range(totalPages):
                page = pdfReader.getPage(i)
                text += page.extractText()
            speechUtil.synthesize_and_save_to_file(text, fileNameWithoutExt)
        elif filetype == 'docx':
            text = docx2txt.process(filePath)
            speechUtil.synthesize_and_save_to_file(text, fileNameWithoutExt)
        elif filetype == 'pptx':
            pptOutput = ''
            ppt = Presentation(filePath)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            speechUtil.synthesize_and_save_to_file(text, fileNameWithoutExt)
    except:
        print("An error occurred!")

    LexicalAnalyzer.perform_lexical_analysis(text)
    return success
        



def clean_up_files():
    files = glob.glob('output/*')
    for f in files:
        os.remove(f)
    files2 = glob.glob('uploads/*')
    for f in files2:
        os.remove(f)

def is_file_empty(filePath):
    with open(filePath, 'r') as fileReader:
        first_char = fileReader.read(1)
        if not first_char:
            return True
    return False